# Copyright 2021, Andres Fredes, <andres.hector.fredes@gmail.com>
# 
# This file is part of pdf2ppt.
# 
#     pdf2ppt is free software: you can redistribute it and/or modify
#     it under the terms of the GNU General Public License as published by
#     the Free Software Foundation, either version 3 of the License, or
#     (at your option) any later version.
# 
#     pdf2ppt is distributed in the hope that it will be useful,
#     but WITHOUT ANY WARRANTY; without even the implied warranty of
#     MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
#     GNU General Public License for more details.
# 
#     You should have received a copy of the GNU General Public License
#     along with pdf2ppt.  If not, see <https://www.gnu.org/licenses/>.

"""coverter.py: Custom worker object to run conversion on secondary thread.
This avoids locking the gui during file processing.
"""
from io import BytesIO
import fitz
from pptx import Presentation
from pptx.util import Cm
from PyQt5.QtCore import QObject, pyqtSignal
from dataclasses import dataclass
from config import BLANK_SLIDE, PPT_TEMPLATE, RATIO_16_9, SLIDE

class Converter(QObject):
    """Simple worker function to handle file conversion in a secondary thread
    to avoid locking the GUI.

    Thread management is done by the GUI itself, but this is the object that
    will be passed to the secondary thread to manage the comparatively slower
    task of converting the files.

    Uses PyMuPDF to turn the pdf's pages into images and python-pptx to put the
    images into a Powerpoint presentation.

    Class variables:
        finished (pyqtSignal): Signal emitted to main thread when task complete.
        warning (pyqtSignal): Signal emitted to main thread when task failed.

    Args:
        path (str): path of file to convert
    """
    finished = pyqtSignal()
    warning = pyqtSignal()

    def __init__(self, path, ppt_cfg):
        super().__init__()
        self.path = path
        self.ppt_cfg = ppt_cfg

    def run(self):
        """Handles the safe file conversion and ultimate emitting of signals.

        This function is triggered when the object is moved onto the secondary
        thread. It emits a finished signal on success and warning on failure.
        """
        try:
            self.convert(self.path, self.ppt_cfg)
        except Exception as e:
            print(e)
            self.warning.emit(e)
        else:
            self.finished.emit()

    def convert(self, path, ppt_cfg):
        """Performs conversion using PyMuPdf and python-pptx

        Args:
            path (str): full system path to pdf file to convert
            ppt_config (Ppt_Config): User configuration variables.
                Defaults to True.
        """
        ppt = Presentation(PPT_TEMPLATE)
        slide_layout = ppt.slide_layouts[BLANK_SLIDE]
        with fitz.open(path) as pdf:
            for page in pdf.pages():
                pixmap = page.getPixmap(
                    matrix=fitz.Matrix(ppt_cfg.img_zoom, ppt_cfg.img_zoom),
                    colorspace=fitz.csRGB,
                    clip=None,
                    alpha=False,
                    annots=ppt_cfg.show_annotations,
                )
                img = {
                    "stream": BytesIO(pixmap.pil_tobytes(
                        format="JPEG",
                        optimize=True,
                    )),
                    "ratio": float(pixmap.w) / float(pixmap.h),
                }
                slide = ppt.slides.add_slide(slide_layout)
                shapes = slide.shapes
                if (img["ratio"] < RATIO_16_9["MIN"]):
                    # Tall ratio, set height to full
                    y = 0
                    height = SLIDE["HEIGHT"]
                    width = height * img["ratio"]
                    x = (SLIDE["WIDTH"] - width) / 2
                elif (img["ratio"] >= RATIO_16_9["MIN"] and
                    img["ratio"] <= RATIO_16_9["MAX"]):
                    # 16:9 ratio, height and width to full
                    x = y = 0
                    width = SLIDE["WIDTH"]
                    height = SLIDE["HEIGHT"]
                else:
                    # Wide ratio, set width to full
                    x = 0
                    width = SLIDE["WIDTH"]
                    height = width / img["ratio"]
                    y = (SLIDE["HEIGHT"] - height) / 2
                shapes.add_picture(
                    img["stream"],
                    Cm(x),
                    Cm(y),
                    Cm(width),
                    Cm(height),
                )
        output_path = f"{path.split('.')[0]}.pptx"
        ppt.save(output_path)


@dataclass
class Ppt_Config():
    """Class for managing the powerpoint configuration variables as selected
    by the user.
    """
    ratio_16_9: bool = True
    page_as_image: bool = True
    img_zoom: int = 2
    show_annotations: bool = False
    
    